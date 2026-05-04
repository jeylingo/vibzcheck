import os
import sys

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Vibzcheck: Collaborative Music App"
    subtitle.text = "Final Project Presentation\nProsper (jeylingo)"

    # 1. Problem and Objectives
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "1. Problem and Objectives"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Problem: The 'Aux Cord Monopoly'"
    p = tf.add_paragraph()
    p.text = "- One person controls the music, leaving others out."
    p = tf.add_paragraph()
    p.text = "- Lack of democratic queue management in standard apps."
    p = tf.add_paragraph()
    p.text = "Core Objectives:"
    p = tf.add_paragraph()
    p.text = "- Inputs: User votes, search queries (iTunes API), chat messages."
    p = tf.add_paragraph()
    p.text = "- Outputs: Live synchronized queue, 30s audio previews, dynamic recommendations."

    # 2. Design and Architecture
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "2. Design and Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "Tech Stack: Flutter & Firebase"
    p = tf.add_paragraph()
    p.text = "- Data Model: Firestore for real-time Rooms, Queues, and Chat."
    p = tf.add_paragraph()
    p.text = "- Architecture: StreamBuilders for reactive UI updates."
    p = tf.add_paragraph()
    p.text = "- UI Flow: Home -> Open Rooms -> Room Queue/Chat."
    p = tf.add_paragraph()
    p.text = "- Aesthetics: Custom Neon Cyberpunk Dark Theme."

    # 3. Testing and Validation
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "3. Testing and Validation"
    tf = slide.placeholders[1].text_frame
    tf.text = "End-to-End User Journey Testing"
    p = tf.add_paragraph()
    p.text = "- Tested room creation and joining via 6-character codes."
    p = tf.add_paragraph()
    p.text = "- Concurrency Testing: Verified Firestore voting transactions."
    p = tf.add_paragraph()
    p.text = "- Playback Sync: Ensured listeners sync within 3 seconds of the host."
    p = tf.add_paragraph()
    p.text = "- Bug Fixed: Bypassed Firebase Storage limits using Base64 string compression."

    # 4. Team and Contributions
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "4. Team and Contributions"
    tf = slide.placeholders[1].text_frame
    tf.text = "Prosper (jeylingo) - Lead Developer"
    p = tf.add_paragraph()
    p.text = "- Product Owner: Vision, UI/UX testing, and overall architecture."
    p = tf.add_paragraph()
    p.text = "- Antigravity (AI): Assisted with API integration and security rules."

    # 5. Presenter Participation
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "5. Presenter Participation"
    tf = slide.placeholders[1].text_frame
    tf.text = "Live Demonstration"
    p = tf.add_paragraph()
    p.text = "- Be prepared to walk through the app live."
    p = tf.add_paragraph()
    p.text = "- Explain the reasoning behind using iTunes API over Spotify."
    p = tf.add_paragraph()
    p.text = "- Demonstrate the live voting and queue reordering."

    # 6. Code Evidence
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "6. Code Evidence & Firebase"
    tf = slide.placeholders[1].text_frame
    tf.text = "Firestore Atomic Transactions (Voting):"
    p = tf.add_paragraph()
    p.text = "await db.runTransaction((transaction) async {\n" \
             "  final songSnap = await transaction.get(songRef);\n" \
             "  int currentScore = songSnap['voteScore'] ?? 0;\n" \
             "  currentScore += voteValue;\n" \
             "  transaction.update(songRef, {'voteScore': currentScore});\n" \
             "});"
    p.font.size = Pt(14)
    p.font.name = "Courier New"

    # Save the presentation
    ppt_path = "Vibzcheck_Presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation generated successfully: {ppt_path}")

if __name__ == "__main__":
    create_presentation()
